from pptx import Presentation
from pptx.util import Pt

# Create a new presentation
prs = Presentation()

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Commercial Banking Intelligence"
subtitle.text = "AI-Driven Future"

# Slide 1: Market Trends
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Market Trends"

body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Increasing adoption of AI in banking operations."

p = tf.add_paragraph()
p.text = "Rising demand for real-time data analytics."
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "Growing competition from fintech startups."
p.font.size = Pt(14)

# Slide 2: Core Technologies
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Core Technologies"

body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Natural Language Processing (NLP)"

p = tf.add_paragraph()
p.text = "Predictive Analytics"
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "Machine Learning Models"
p.font.size = Pt(14)

# Slide 3: Strategic Benefits
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Strategic Benefits"

body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Enhanced decision-making capabilities."

p = tf.add_paragraph()
p.text = "Improved customer experience through personalization."
p.font.size = Pt(14)

p = tf.add_paragraph()
p.text = "Cost reduction via automation and efficiency gains."
p.font.size = Pt(14)

# Save the presentation
prs.save('commercial_banking_ai.pptx')